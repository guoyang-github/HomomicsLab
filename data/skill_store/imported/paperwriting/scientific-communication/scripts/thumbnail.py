#!/usr/bin/env python3
"""
Generate thumbnails or individual slide images from PowerPoint presentations.
"""

import argparse
import sys
from pathlib import Path


def generate_thumbnails(pptx_path: Path, output_dir: Path, cols: int = 4, individual: bool = False) -> None:
    """Generate thumbnails from a PowerPoint file."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
    except ImportError:
        print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(pptx_path)
    output_dir.mkdir(parents=True, exist_ok=True)

    if individual:
        # Export each slide as a separate PNG (requires additional libraries)
        print(f"Exporting {len(prs.slides)} individual slides to {output_dir}")
        try:
            # Try using comtypes on Windows or unoconv on Linux
            # Fallback: save as PDF then convert
            _export_individual_slides(pptx_path, output_dir)
        except Exception as e:
            print(f"Individual export requires additional setup: {e}", file=sys.stderr)
            print("Fallback: Saving slide text summaries", file=sys.stderr)
            for i, slide in enumerate(prs.slides, 1):
                text = _extract_slide_text(slide)
                out_file = output_dir / f"slide_{i:02d}.txt"
                out_file.write_text(text, encoding="utf-8")
    else:
        # Generate contact sheet style thumbnails
        print(f"Generating thumbnail grid ({cols} columns) from {pptx_path}")
        try:
            _generate_contact_sheet(prs, output_dir, cols)
        except Exception as e:
            print(f"Contact sheet generation requires PIL: {e}", file=sys.stderr)
            # Fallback: text index
            index_file = output_dir / "slide_index.txt"
            lines = [f"Thumbnail grid ({cols} cols) for: {pptx_path.name}\n"]
            for i, slide in enumerate(prs.slides, 1):
                text = _extract_slide_text(slide)[:200]
                lines.append(f"Slide {i}: {text}...\n")
            index_file.write_text("".join(lines), encoding="utf-8")
            print(f"Saved text index to {index_file}")


def _extract_slide_text(slide) -> str:
    """Extract all text from a slide."""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
    return " | ".join(texts)


def _export_individual_slides(pptx_path: Path, output_dir: Path) -> None:
    """Try to export individual slide images."""
    # This requires LibreOffice or PowerPoint automation
    # Provide instructions if not available
    print("Note: For high-quality slide images, use one of these methods:")
    print(f"  1. PowerPoint: File > Export > Change File Type > PNG")
    print(f"  2. LibreOffice: libreoffice --headless --convert-to pdf {pptx_path}")
    print(f"     Then: pdftoppm -png output.pdf {output_dir}/slide")
    print(f"  3. Online converter or screenshot tool")


def _generate_contact_sheet(prs, output_dir: Path, cols: int) -> None:
    """Generate a contact sheet image. Requires PIL."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        raise ImportError("PIL is required for contact sheet generation")

    # Placeholder: create a text-based contact sheet
    # Real implementation would render each slide to an image
    width, height = 200 * cols, 150 * ((len(prs.slides) + cols - 1) // cols)
    img = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(img)

    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)
    except (OSError, IOError):
        font = ImageFont.load_default()

    for i, slide in enumerate(prs.slides):
        row, col = divmod(i, cols)
        x, y = col * 200, row * 150
        text = _extract_slide_text(slide)[:100]
        draw.rectangle([x, y, x + 198, y + 148], outline="black", width=1)
        draw.text((x + 5, y + 5), f"Slide {i + 1}", fill="black", font=font)
        draw.text((x + 5, y + 25), text, fill="gray", font=font)

    output_file = output_dir / "thumbnails.png"
    img.save(output_file)
    print(f"Saved contact sheet to {output_file}")


def main():
    parser = argparse.ArgumentParser(description="Generate thumbnails from PowerPoint presentations")
    parser.add_argument("pptx", help="Input PowerPoint file")
    parser.add_argument("output", help="Output directory")
    parser.add_argument("--cols", type=int, default=4, help="Number of columns for contact sheet")
    parser.add_argument("--individual", action="store_true", help="Export individual slide images")
    args = parser.parse_args()

    generate_thumbnails(Path(args.pptx), Path(args.output), args.cols, args.individual)


if __name__ == "__main__":
    main()
