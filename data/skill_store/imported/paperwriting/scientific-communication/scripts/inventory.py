#!/usr/bin/env python3
"""
Extract text inventory from PowerPoint presentations to JSON.
Useful for content review, translation workflows, and bulk replacements.
"""

import argparse
import json
import sys
from pathlib import Path
from typing import Dict, List


def extract_inventory(pptx_path: Path) -> Dict:
    """Extract all text content from a PowerPoint file."""
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(pptx_path)
    slides_data = []

    for i, slide in enumerate(prs.slides, 1):
        slide_info = {
            "slide_number": i,
            "layout": slide.slide_layout.name if slide.slide_layout else "unknown",
            "shapes": []
        }

        for shape in slide.shapes:
            shape_info = {
                "shape_id": shape.shape_id,
                "shape_type": str(shape.shape_type),
                "name": shape.name,
                "text": "",
                "has_text_frame": False
            }

            if shape.has_text_frame:
                shape_info["has_text_frame"] = True
                texts = []
                for paragraph in shape.text_frame.paragraphs:
                    para_text = "".join(run.text for run in paragraph.runs)
                    if para_text.strip():
                        texts.append(para_text)
                shape_info["text"] = "\n".join(texts)

            slide_info["shapes"].append(shape_info)

        slides_data.append(slide_info)

    return {
        "source_file": str(pptx_path),
        "total_slides": len(prs.slides),
        "slides": slides_data
    }


def save_inventory(inventory: Dict, output_path: Path) -> None:
    """Save inventory to JSON file."""
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(inventory, f, indent=2, ensure_ascii=False)
    print(f"Inventory saved to {output_path}")
    print(f"Total slides: {inventory['total_slides']}")

    # Print summary
    total_shapes = sum(len(s["shapes"]) for s in inventory["slides"])
    text_shapes = sum(1 for s in inventory["slides"] for sh in s["shapes"] if sh["text"])
    print(f"Total shapes: {total_shapes}")
    print(f"Shapes with text: {text_shapes}")


def main():
    parser = argparse.ArgumentParser(description="Extract text inventory from PowerPoint presentations")
    parser.add_argument("pptx", help="Input PowerPoint file")
    parser.add_argument("output", help="Output JSON file")
    args = parser.parse_args()

    inventory = extract_inventory(Path(args.pptx))
    save_inventory(inventory, Path(args.output))


if __name__ == "__main__":
    main()
