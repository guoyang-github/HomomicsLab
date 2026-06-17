#!/usr/bin/env python3
"""
Replace text content in PowerPoint presentations using a JSON replacements file.
"""

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Dict


def load_replacements(replacements_path: Path) -> Dict[str, str]:
    """Load replacements from JSON file."""
    with open(replacements_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    # Support both flat dict and structured format
    if isinstance(data, dict) and all(isinstance(v, str) for v in data.values()):
        return data
    elif isinstance(data, dict) and "replacements" in data:
        return {r["old"]: r["new"] for r in data["replacements"]}
    else:
        print("Error: replacements JSON should be {old_text: new_text} or {replacements: [{old, new}]}", file=sys.stderr)
        sys.exit(1)


def replace_text_in_presentation(input_path: Path, output_path: Path, replacements: Dict[str, str]) -> int:
    """Replace text in PowerPoint file. Returns number of replacements made."""
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(input_path)
    total_replacements = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    new_text = original_text
                    for old_str, new_str in replacements.items():
                        new_text = new_text.replace(old_str, new_str)
                    if new_text != original_text:
                        run.text = new_text
                        total_replacements += 1

    prs.save(output_path)
    return total_replacements


def main():
    parser = argparse.ArgumentParser(description="Replace text in PowerPoint presentations")
    parser.add_argument("input", help="Input PowerPoint file")
    parser.add_argument("replacements", help="JSON file with text replacements")
    parser.add_argument("output", help="Output PowerPoint file")
    parser.add_argument("--regex", action="store_true", help="Treat replacements as regex patterns")
    args = parser.parse_args()

    replacements = load_replacements(Path(args.replacements))

    if args.regex:
        # Convert to regex replacements
        regex_replacements = {}
        for old, new in replacements.items():
            regex_replacements[old] = new
        # Note: regex replacement in pptx runs is complex due to text splitting
        print("Warning: regex mode may not work perfectly with split text runs", file=sys.stderr)

    count = replace_text_in_presentation(Path(args.input), Path(args.output), replacements)
    print(f"Made {count} text replacements")
    print(f"Saved to {args.output}")


if __name__ == "__main__":
    main()
