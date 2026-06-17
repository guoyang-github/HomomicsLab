"""
Academic PPTX backend for scientific presentation generation.

Generates structured PowerPoint slides optimized for academic and scientific
contexts. Supports multiple layouts, themes, chart embedding, and image
insertion — all using only python-pptx (zero external API required).

Layouts: title, content, two_column, image_text, chart, table, full_image
Themes: light (default), dark, minimal

Usage:
    python generate_slide_image.py "[LAYOUT: title] My Talk" -o title.pptx
    python generate_slide_image.py "[LAYOUT: chart] Results\nDATA: x=[1,2,3], y=[4,5,6]" -o results.pptx
    python generate_slide_image.py "[LAYOUT: image_text] Methods\nIMAGE: figures/method.png" -o methods.pptx
"""

import os
import re
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from . import Backend


# ── Academic Themes ────────────────────────────────────────────────────
# Okabe-Ito colorblind-friendly palette used as accents
OKABE_ITO = {
    'orange': (230, 159, 0),      # #E69F00
    'sky_blue': (86, 180, 233),   # #56B4E9
    'green': (0, 158, 115),       # #009E73
    'yellow': (240, 228, 66),     # #F0E442
    'blue': (0, 114, 178),        # #0072B2
    'vermillion': (213, 94, 0),   # #D55E00
    'purple': (204, 121, 167),    # #CC79A7
    'black': (0, 0, 0),
}

THEMES = {
    'light': {
        'name': 'light',
        'bg': (255, 255, 255),
        'title': (33, 33, 33),
        'subtitle': (80, 80, 80),
        'text': (68, 68, 68),
        'bullet': (68, 68, 68),
        'accent': OKABE_ITO['blue'],
        'secondary': OKABE_ITO['orange'],
        'font_title': 'Calibri',
        'font_body': 'Calibri',
        'font_mono': 'Courier New',
    },
    'dark': {
        'name': 'dark',
        'bg': (26, 35, 126),          # deep blue
        'title': (255, 255, 255),
        'subtitle': (200, 200, 200),
        'text': (220, 220, 220),
        'bullet': (220, 220, 220),
        'accent': (255, 193, 7),      # gold
        'secondary': (144, 202, 249), # light blue
        'font_title': 'Calibri',
        'font_body': 'Calibri',
        'font_mono': 'Consolas',
    },
    'minimal': {
        'name': 'minimal',
        'bg': (250, 250, 250),
        'title': (0, 0, 0),
        'subtitle': (80, 80, 80),
        'text': (80, 80, 80),
        'bullet': (80, 80, 80),
        'accent': OKABE_ITO['orange'],
        'secondary': OKABE_ITO['blue'],
        'font_title': 'Helvetica',
        'font_body': 'Arial',
        'font_mono': 'Courier New',
    },
}

# ── Layout Capacity Matrix (inspired by Mck-ppt-design) ────────────────
# Char budgets prevent overflow — key for programmatic slide generation
LAYOUT_MATRIX = {
    'title':      {'title_chars': 80,  'subtitle_chars': 120, 'author_chars': 60},
    'content':    {'title_chars': 80,  'bullet_chars': 100,   'max_bullets': 6},
    'two_column': {'title_chars': 80,  'bullet_chars': 80,    'max_bullets': 5, 'max_per_col': 5},
    'image_text': {'title_chars': 80,  'caption_chars': 200,  'text_chars': 400},
    'chart':      {'title_chars': 80,  'caption_chars': 200},
    'table':      {'title_chars': 80,  'cell_chars': 60,      'max_rows': 8, 'max_cols': 5},
    'full_image': {'title_chars': 80,  'caption_chars': 200},
}

SLIDE_W_INCH = 13.333
SLIDE_H_INCH = 7.5


class PPTXBackend(Backend):
    """Generate academic/scientific slides using python-pptx with multi-layout support."""

    name = "pptx"

    def __init__(self, script_dir: Path):
        self.script_dir = script_dir
        try:
            import pptx
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE
            self._pptx = pptx
            self._Inches = Inches
            self._Pt = Pt
            self._Emu = Emu
            self._RGBColor = RGBColor
            self._PP_ALIGN = PP_ALIGN
            self._MSO_ANCHOR = MSO_ANCHOR
            self._MSO_SHAPE = MSO_SHAPE
        except ImportError:
            raise ImportError(
                "python-pptx required for PPTX backend. Install: pip install python-pptx"
            )

    # ═══════════════════════════════════════════════════════════════════
    # Public API
    # ═══════════════════════════════════════════════════════════════════

    def generate(self, prompt: str, output_path: Path, visual_only: bool = False,
                 iterations: int = 2, verbose: bool = False) -> Dict[str, Any]:
        if iterations > 1 and verbose:
            print(f"[PPTX] Note: iterations={iterations} is ignored by the pptx backend")

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Detect theme
        theme_name = self._detect_theme(prompt)
        theme = THEMES.get(theme_name, THEMES['light'])

        if visual_only:
            return self._generate_visual(prompt, output_path, theme, verbose)

        # Full slide
        if output_path.suffix.lower() != ".pptx":
            output_path = output_path.with_suffix(".pptx")

        layout = self._detect_layout(prompt, visual_only=False)
        prs = self._create_presentation(theme)

        # Route to layout handler
        handler = getattr(self, f'_layout_{layout}', self._layout_content)
        try:
            handler(prs, prompt, theme, verbose)
        except Exception as e:
            if verbose:
                print(f"[PPTX] Layout '{layout}' failed: {e}, falling back to content")
            self._layout_content(prs, prompt, theme, verbose)

        # QA validation
        qa_issues = self._qa_validate(prs)
        if qa_issues and verbose:
            print(f"[PPTX] QA: {len(qa_issues)} issues found")
            for issue in qa_issues:
                print(f"  - {issue}")

        prs.save(str(output_path))

        return {
            "success": True,
            "final_image": str(output_path),
            "mode": "pptx",
            "layout": layout,
            "theme": theme_name,
            "qa_issues": len(qa_issues),
            "iterations_used": 1
        }

    # ═══════════════════════════════════════════════════════════════════
    # Theme & Layout Detection
    # ═══════════════════════════════════════════════════════════════════

    def _detect_theme(self, prompt: str) -> str:
        """Detect theme from prompt keywords."""
        p = prompt.lower()
        if '[theme: dark]' in p or 'dark background' in p or 'dark mode' in p:
            return 'dark'
        if '[theme: minimal]' in p or 'minimal' in p:
            return 'minimal'
        if '[theme: light]' in p:
            return 'light'
        return 'light'

    def _detect_layout(self, prompt: str, visual_only: bool) -> str:
        """Detect layout type from explicit tags or content inference."""
        p = prompt.lower()

        # Explicit layout tag
        m = re.search(r'\[LAYOUT:\s*(\w+)\]', prompt, re.IGNORECASE)
        if m:
            detected = m.group(1).lower()
            if detected in LAYOUT_MATRIX:
                return detected
            # Aliases
            aliases = {
                'title': 'title', 'cover': 'title', 'heading': 'title',
                'content': 'content', 'bullets': 'content', 'list': 'content',
                'two_column': 'two_column', 'twocolumn': 'two_column', 'compare': 'two_column',
                'image_text': 'image_text', 'image': 'image_text', 'figure_text': 'image_text',
                'chart': 'chart', 'graph': 'chart', 'plot': 'chart',
                'table': 'table', 'data_table': 'table',
                'full_image': 'full_image', 'fullimage': 'full_image', 'full_figure': 'full_image',
            }
            if detected in aliases:
                return aliases[detected]

        # Visual-only inference
        if visual_only:
            if re.search(r'\bDATA:\s*x\s*=\s*\[', prompt, re.IGNORECASE):
                return 'chart'
            if re.search(r'\bIMAGE:\s*\S+\.(?:png|jpg|jpeg)', prompt, re.IGNORECASE):
                return 'full_image'
            return 'content'

        # Full-slide inference
        if re.search(r'^(?:title|cover)\b', prompt, re.IGNORECASE):
            return 'title'
        if re.search(r'\bDATA:\s*x\s*=\s*\[', prompt, re.IGNORECASE):
            return 'chart'
        if re.search(r'\bTABLE:\s*', prompt, re.IGNORECASE) or (
                'table' in p and '|' in prompt):
            return 'table'
        if re.search(r'\bIMAGE:\s*\S+\.(?:png|jpg|jpeg)', prompt, re.IGNORECASE):
            return 'image_text'
        if 'two column' in p or 'left and right' in p or 'compare' in p or 'vs' in p:
            return 'two_column'
        if 'full image' in p or 'full figure' in p:
            return 'full_image'
        return 'content'

    # ═══════════════════════════════════════════════════════════════════
    # Presentation Factory
    # ═══════════════════════════════════════════════════════════════════

    def _create_presentation(self, theme: Dict[str, Any]):
        """Create a new presentation with theme settings."""
        prs = self._pptx.Presentation()
        prs.slide_width = self._Inches(SLIDE_W_INCH)
        prs.slide_height = self._Inches(SLIDE_H_INCH)
        return prs

    def _blank_slide(self, prs):
        """Add a blank slide (layout 6 is blank in most templates)."""
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _rgb(self, rgb_tuple):
        """Helper to create RGBColor."""
        return self._RGBColor(*rgb_tuple)

    # ═══════════════════════════════════════════════════════════════════
    # Layout Handlers — Full Slides
    # ═══════════════════════════════════════════════════════════════════

    def _layout_title(self, prs, prompt: str, theme: Dict[str, Any], verbose: bool):
        """Title/cover slide: large title, subtitle, author, date."""
        s = self._blank_slide(prs)

        # Background
        bg = s.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE, 0, 0,
            self._Inches(SLIDE_W_INCH), self._Inches(SLIDE_H_INCH)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(theme['bg'])
        bg.line.fill.background()

        # Accent bar at top
        bar_h = self._Inches(0.08)
        bar = s.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE, 0, 0,
            self._Inches(SLIDE_W_INCH), bar_h
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(theme['accent'])
        bar.line.fill.background()

        # Parse fields
        title = self._extract_field(prompt, ['title', 'Title'])
        subtitle = self._extract_field(prompt, ['subtitle', 'Subtitle'])
        author = self._extract_field(prompt, ['author', 'Author', 'speaker', 'Speaker'])
        date = self._extract_field(prompt, ['date', 'Date'])

        # Fallback: first line as title
        if not title:
            first_line = prompt.split('\n')[0].strip()
            first_line = re.sub(r'\[LAYOUT:\s*\w+\]', '', first_line, flags=re.IGNORECASE).strip()
            title = first_line[:LAYOUT_MATRIX['title']['title_chars']]

        # Title (large, centered)
        title_box = s.shapes.add_textbox(
            self._Inches(1), self._Inches(2.2),
            self._Inches(SLIDE_W_INCH - 2), self._Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self._Pt(54)
        p.font.bold = True
        p.font.name = theme['font_title']
        p.font.color.rgb = self._rgb(theme['title'])
        p.alignment = self._PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            sub_box = s.shapes.add_textbox(
                self._Inches(1), self._Inches(3.5),
                self._Inches(SLIDE_W_INCH - 2), self._Inches(0.8)
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = self._Pt(28)
            p.font.name = theme['font_body']
            p.font.color.rgb = self._rgb(theme['subtitle'])
            p.alignment = self._PP_ALIGN.CENTER

        # Author + Date (bottom)
        info_parts = [p for p in [author, date] if p]
        if info_parts:
            info_box = s.shapes.add_textbox(
                self._Inches(1), self._Inches(5.8),
                self._Inches(SLIDE_W_INCH - 2), self._Inches(0.4)
            )
            tf = info_box.text_frame
            p = tf.paragraphs[0]
            p.text = '  |  '.join(info_parts)
            p.font.size = self._Pt(18)
            p.font.name = theme['font_body']
            p.font.color.rgb = self._rgb(theme['subtitle'])
            p.alignment = self._PP_ALIGN.CENTER

    def _layout_content(self, prs, prompt: str, theme: Dict[str, Any], verbose: bool):
        """Standard content slide: title + bullet points."""
        s = self._blank_slide(prs)

        # Background
        bg = s.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE, 0, 0,
            self._Inches(SLIDE_W_INCH), self._Inches(SLIDE_H_INCH)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(theme['bg'])
        bg.line.fill.background()

        parsed = self._parse_prompt(prompt)
        title = parsed["title"]
        bullets = parsed["bullets"]
        limits = LAYOUT_MATRIX['content']

        # Truncate to limits
        title = title[:limits['title_chars']]
        bullets = bullets[:limits['max_bullets']]
        bullets = [b[:limits['bullet_chars']] for b in bullets]

        # Accent bar
        bar = s.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE, 0, 0,
            self._Inches(SLIDE_W_INCH), self._Inches(0.06)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(theme['accent'])
        bar.line.fill.background()

        # Title
        title_box = s.shapes.add_textbox(
            self._Inches(0.6), self._Inches(0.4),
            self._Inches(SLIDE_W_INCH - 1.2), self._Inches(0.8)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self._Pt(40)
        p.font.bold = True
        p.font.name = theme['font_title']
        p.font.color.rgb = self._rgb(theme['title'])

        # Bullets
        if bullets:
            content_box = s.shapes.add_textbox(
                self._Inches(0.6), self._Inches(1.4),
                self._Inches(SLIDE_W_INCH - 1.2), self._Inches(5.0)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = self._Pt(22)
                p.font.name = theme['font_body']
                p.font.color.rgb = self._rgb(theme['bullet'])
                p.space_after = self._Pt(16)

    def _layout_two_column(self, prs, prompt: str, theme: Dict[str, Any], verbose: bool):
        """Two-column comparison slide."""
        s = self._blank_slide(prs)

        # Background
        bg = s.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE, 0, 0,
            self._Inches(SLIDE_W_INCH), self._Inches(SLIDE_H_INCH)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(theme['bg'])
        bg.line.fill.background()

        parsed = self._parse_prompt(prompt)
        title = parsed["title"][:LAYOUT_MATRIX['two_column']['title_chars']]

        # Try to split bullets into left/right
        bullets = parsed["bullets"]
        mid = len(bullets) // 2
        left_bullets = bullets[:mid][:LAYOUT_MATRIX['two_column']['max_per_col']]
        right_bullets = bullets[mid:][:LAYOUT_MATRIX['two_column']['max_per_col']]

        # Accent bar
        bar = s.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE, 0, 0,
            self._Inches(SLIDE_W_INCH), self._Inches(0.06)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(theme['accent'])
        bar.line.fill.background()

        # Title
        title_box = s.shapes.add_textbox(
            self._Inches(0.6), self._Inches(0.4),
            self._Inches(SLIDE_W_INCH - 1.2), self._Inches(0.7)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self._Pt(36)
        p.font.bold = True
        p.font.name = theme['font_title']
        p.font.color.rgb = self._rgb(theme['title'])

        # Divider line
        line = s.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE,
            self._Inches(SLIDE_W_INCH / 2 - 0.01), self._Inches(1.3),
            self._Inches(0.02), self._Inches(5.0)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._rgb(theme['accent'])
        line.line.fill.background()

        # Left column
        self._add_bullet_column(s, left_bullets, theme,
                                self._Inches(0.6), self._Inches(1.4),
                                self._Inches(5.5), self._Inches(5.0))

        # Right column
        self._add_bullet_column(s, right_bullets, theme,
                                self._Inches(SLIDE_W_INCH / 2 + 0.3), self._Inches(1.4),
                                self._Inches(5.5), self._Inches(5.0))

    def _layout_image_text(self, prs, prompt: str, theme: Dict[str, Any], verbose: bool):
        """Image + text slide. Detects IMAGE: path in prompt."""
        s = self._blank_slide(prs)

        # Background
        bg = s.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE, 0, 0,
            self._Inches(SLIDE_W_INCH), self._Inches(SLIDE_H_INCH)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(theme['bg'])
        bg.line.fill.background()

        parsed = self._parse_prompt(prompt)
        title = parsed["title"][:LAYOUT_MATRIX['image_text']['title_chars']]
        image_paths = self._parse_image_paths(prompt)

        # Accent bar
        bar = s.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE, 0, 0,
            self._Inches(SLIDE_W_INCH), self._Inches(0.06)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(theme['accent'])
        bar.line.fill.background()

        # Title
        title_box = s.shapes.add_textbox(
            self._Inches(0.6), self._Inches(0.4),
            self._Inches(SLIDE_W_INCH - 1.2), self._Inches(0.7)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self._Pt(36)
        p.font.bold = True
        p.font.name = theme['font_title']
        p.font.color.rgb = self._rgb(theme['title'])

        # Try to place image on left, text on right
        img_left = self._Inches(0.6)
        img_top = self._Inches(1.3)
        img_w = self._Inches(5.5)
        img_h = self._Inches(4.8)
        text_left = self._Inches(6.8)
        text_w = self._Inches(5.8)

        if image_paths:
            img_path = image_paths[0]
            if os.path.isfile(img_path):
                try:
                    s.shapes.add_picture(img_path, img_left, img_top, width=img_w)
                except Exception as e:
                    if verbose:
                        print(f"[PPTX] Could not load image {img_path}: {e}")
                    self._add_placeholder_rect(s, img_left, img_top, img_w, img_h, theme)
            else:
                if verbose:
                    print(f"[PPTX] Image not found: {img_path}")
                self._add_placeholder_rect(s, img_left, img_top, img_w, img_h, theme)
        else:
            self._add_placeholder_rect(s, img_left, img_top, img_w, img_h, theme)

        # Text on right
        text_box = s.shapes.add_textbox(text_left, img_top, text_w, img_h)
        tf = text_box.text_frame
        tf.word_wrap = True
        bullets = parsed["bullets"][:5]
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet[:LAYOUT_MATRIX['image_text']['text_chars']]}"

            p.font.size = self._Pt(20)
            p.font.name = theme['font_body']
            p.font.color.rgb = self._rgb(theme['text'])
            p.space_after = self._Pt(12)

    def _layout_chart(self, prs, prompt: str, theme: Dict[str, Any], verbose: bool):
        """Chart slide: generates matplotlib chart from DATA: marker."""
        s = self._blank_slide(prs)

        # Background
        bg = s.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE, 0, 0,
            self._Inches(SLIDE_W_INCH), self._Inches(SLIDE_H_INCH)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(theme['bg'])
        bg.line.fill.background()

        parsed = self._parse_prompt(prompt)
        title = parsed["title"][:LAYOUT_MATRIX['chart']['title_chars']]

        # Accent bar
        bar = s.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE, 0, 0,
            self._Inches(SLIDE_W_INCH), self._Inches(0.06)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(theme['accent'])
        bar.line.fill.background()

        # Title
        title_box = s.shapes.add_textbox(
            self._Inches(0.6), self._Inches(0.4),
            self._Inches(SLIDE_W_INCH - 1.2), self._Inches(0.7)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self._Pt(36)
        p.font.bold = True
        p.font.name = theme['font_title']
        p.font.color.rgb = self._rgb(theme['title'])

        # Generate chart
        chart_data = self._parse_chart_data(prompt)
        if chart_data:
            chart_path = self._generate_matplotlib_chart(chart_data, theme, verbose)
            if chart_path and os.path.isfile(chart_path):
                s.shapes.add_picture(
                    chart_path,
                    self._Inches(1.5), self._Inches(1.3),
                    width=self._Inches(10.0)
                )
                try:
                    os.remove(chart_path)
                except OSError:
                    pass
        else:
            # No chart data: show placeholder with instructions
            placeholder = s.shapes.add_textbox(
                self._Inches(2), self._Inches(2.5),
                self._Inches(9), self._Inches(2)
            )
            tf = placeholder.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = (
                "To generate a chart, include data in your prompt:\n"
                "DATA: x=[1,2,3,4,5], y=[10,20,15,30,25]"
            )
            p.font.size = self._Pt(18)
            p.font.name = theme['font_body']
            p.font.color.rgb = self._rgb(theme['subtitle'])

        # Caption
        caption = self._extract_field(prompt, ['caption', 'Caption', 'note', 'Note'])
        if caption:
            cap_box = s.shapes.add_textbox(
                self._Inches(0.6), self._Inches(6.5),
                self._Inches(SLIDE_W_INCH - 1.2), self._Inches(0.5)
            )
            tf = cap_box.text_frame
            p = tf.paragraphs[0]
            p.text = caption[:LAYOUT_MATRIX['chart']['caption_chars']]
            p.font.size = self._Pt(14)
            p.font.italic = True
            p.font.name = theme['font_body']
            p.font.color.rgb = self._rgb(theme['subtitle'])

    def _layout_table(self, prs, prompt: str, theme: Dict[str, Any], verbose: bool):
        """Table slide: parses TABLE: marker or markdown-style table."""
        s = self._blank_slide(prs)

        # Background
        bg = s.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE, 0, 0,
            self._Inches(SLIDE_W_INCH), self._Inches(SLIDE_H_INCH)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(theme['bg'])
        bg.line.fill.background()

        parsed = self._parse_prompt(prompt)
        title = parsed["title"][:LAYOUT_MATRIX['table']['title_chars']]
        table_data = self._parse_table_data(prompt)

        # Accent bar
        bar = s.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE, 0, 0,
            self._Inches(SLIDE_W_INCH), self._Inches(0.06)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(theme['accent'])
        bar.line.fill.background()

        # Title
        title_box = s.shapes.add_textbox(
            self._Inches(0.6), self._Inches(0.4),
            self._Inches(SLIDE_W_INCH - 1.2), self._Inches(0.7)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self._Pt(36)
        p.font.bold = True
        p.font.name = theme['font_title']
        p.font.color.rgb = self._rgb(theme['title'])

        if table_data and table_data['headers'] and table_data['rows']:
            headers = table_data['headers'][:LAYOUT_MATRIX['table']['max_cols']]
            rows = table_data['rows'][:LAYOUT_MATRIX['table']['max_rows']]
            n_cols = len(headers)
            n_rows = len(rows)

            table_shape = s.shapes.add_table(
                n_rows + 1, n_cols,
                self._Inches(0.8), self._Inches(1.4),
                self._Inches(SLIDE_W_INCH - 1.6), self._Inches(4.8)
            ).table

            # Header row
            for i, h in enumerate(headers):
                cell = table_shape.cell(0, i)
                cell.text = str(h)[:LAYOUT_MATRIX['table']['cell_chars']]
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._rgb(theme['accent'])
                p = cell.text_frame.paragraphs[0]
                p.font.size = self._Pt(16)
                p.font.bold = True
                p.font.name = theme['font_body']
                p.font.color.rgb = self._RGBColor(255, 255, 255)

            # Data rows
            for r_idx, row in enumerate(rows):
                for c_idx in range(n_cols):
                    cell = table_shape.cell(r_idx + 1, c_idx)
                    val = str(row[c_idx]) if c_idx < len(row) else ''
                    cell.text = val[:LAYOUT_MATRIX['table']['cell_chars']]
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = self._Pt(14)
                    p.font.name = theme['font_body']
                    p.font.color.rgb = self._rgb(theme['text'])
                    if r_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self._RGBColor(245, 245, 245)
        else:
            placeholder = s.shapes.add_textbox(
                self._Inches(2), self._Inches(2.5),
                self._Inches(9), self._Inches(2)
            )
            tf = placeholder.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = (
                "To generate a table, use:\n"
                "TABLE: Header1 | Header2 | Header3\n"
                "Row1A | Row1B | Row1C\n"
                "Row2A | Row2B | Row2C"
            )
            p.font.size = self._Pt(18)
            p.font.name = theme['font_body']
            p.font.color.rgb = self._rgb(theme['subtitle'])

    def _layout_full_image(self, prs, prompt: str, theme: Dict[str, Any], verbose: bool):
        """Full-image slide with optional caption."""
        s = self._blank_slide(prs)

        image_paths = self._parse_image_paths(prompt)

        if image_paths and os.path.isfile(image_paths[0]):
            try:
                s.shapes.add_picture(
                    image_paths[0],
                    self._Inches(0.5), self._Inches(0.5),
                    width=self._Inches(SLIDE_W_INCH - 1.0)
                )
            except Exception as e:
                if verbose:
                    print(f"[PPTX] Could not load image: {e}")
                self._add_placeholder_rect(
                    s, self._Inches(0.5), self._Inches(0.5),
                    self._Inches(SLIDE_W_INCH - 1.0), self._Inches(6.0), theme
                )
        else:
            self._add_placeholder_rect(
                s, self._Inches(0.5), self._Inches(0.5),
                self._Inches(SLIDE_W_INCH - 1.0), self._Inches(6.0), theme
            )

        # Caption
        caption = self._extract_field(prompt, ['caption', 'Caption'])
        if not caption:
            # Use bullets as caption
            parsed = self._parse_prompt(prompt)
            if parsed['bullets']:
                caption = parsed['bullets'][0][:LAYOUT_MATRIX['full_image']['caption_chars']]

        if caption:
            cap_box = s.shapes.add_textbox(
                self._Inches(0.5), self._Inches(6.6),
                self._Inches(SLIDE_W_INCH - 1.0), self._Inches(0.5)
            )
            tf = cap_box.text_frame
            p = tf.paragraphs[0]
            p.text = caption
            p.font.size = self._Pt(14)
            p.font.italic = True
            p.font.name = theme['font_body']
            p.font.color.rgb = self._rgb(theme['subtitle'])

    # ═══════════════════════════════════════════════════════════════════
    # Visual-Only Generation
    # ═══════════════════════════════════════════════════════════════════

    def _generate_visual(self, prompt: str, output_path: Path,
                         theme: Dict[str, Any], verbose: bool) -> Dict[str, Any]:
        """Generate a visual-only image (chart or diagram)."""
        output_path = Path(output_path)
        if not output_path.suffix:
            output_path = output_path.with_suffix(".png")

        chart_data = self._parse_chart_data(prompt)
        if chart_data:
            chart_path = self._generate_matplotlib_chart(chart_data, theme, verbose)
            if chart_path and os.path.isfile(chart_path):
                import shutil
                shutil.move(chart_path, str(output_path))
                return {
                    "success": True,
                    "final_image": str(output_path),
                    "mode": "pptx",
                    "visual_type": "chart",
                    "iterations_used": 1
                }

        image_paths = self._parse_image_paths(prompt)
        if image_paths and os.path.isfile(image_paths[0]):
            import shutil
            shutil.copy(image_paths[0], str(output_path))
            return {
                "success": True,
                "final_image": str(output_path),
                "mode": "pptx",
                "visual_type": "image",
                "iterations_used": 1
            }

        # Fallback: matplotlib diagram (enhanced version)
        success = self._create_diagram_visual(output_path, prompt, theme, verbose)
        if success:
            return {
                "success": True,
                "final_image": str(output_path),
                "mode": "pptx",
                "visual_type": "diagram",
                "iterations_used": 1
            }
        return {
            "success": False,
            "error": "Failed to create visual",
            "mode": "pptx"
        }

    # ═══════════════════════════════════════════════════════════════════
    # Chart & Diagram Generation
    # ═══════════════════════════════════════════════════════════════════

    def _parse_chart_data(self, prompt: str) -> Optional[Dict[str, Any]]:
        """Parse DATA: marker from prompt."""
        # DATA: x=[1,2,3], y=[4,5,6]
        m = re.search(
            r'DATA:\s*x\s*=\s*\[([^\]]+)\]\s*,\s*y\s*=\s*\[([^\]]+)\]',
            prompt, re.IGNORECASE
        )
        if m:
            x_vals = [float(v.strip()) for v in m.group(1).split(',') if v.strip()]
            y_vals = [float(v.strip()) for v in m.group(2).split(',') if v.strip()]
            chart_type = 'line'
            if 'bar' in prompt.lower():
                chart_type = 'bar'
            elif 'scatter' in prompt.lower():
                chart_type = 'scatter'
            return {'x': x_vals, 'y': y_vals, 'type': chart_type}
        return None

    def _generate_matplotlib_chart(self, chart_data: Dict[str, Any],
                                   theme: Dict[str, Any], verbose: bool) -> Optional[str]:
        """Generate a matplotlib chart and save to temp file."""
        try:
            import matplotlib
            matplotlib.use('Agg')
            import matplotlib.pyplot as plt

            fig, ax = plt.subplots(figsize=(10, 5.5))

            # Apply theme colors
            bg_rgb = theme['bg']
            ax.set_facecolor([c / 255.0 for c in bg_rgb])
            fig.patch.set_facecolor([c / 255.0 for c in bg_rgb])
            ax.tick_params(colors=[c / 255.0 for c in theme['text']])
            ax.xaxis.label.set_color([c / 255.0 for c in theme['text']])
            ax.yaxis.label.set_color([c / 255.0 for c in theme['text']])
            ax.title.set_color([c / 255.0 for c in theme['title']])
            for spine in ax.spines.values():
                spine.set_color([c / 255.0 for c in theme['subtitle']])

            x = chart_data['x']
            y = chart_data['y']
            accent = [c / 255.0 for c in theme['accent']]
            secondary = [c / 255.0 for c in theme.get('secondary', theme['accent'])]

            if chart_data['type'] == 'bar':
                ax.bar(range(len(x)), y, color=accent, edgecolor='white', linewidth=0.5)
                ax.set_xticks(range(len(x)))
                ax.set_xticklabels([str(v) for v in x])
            elif chart_data['type'] == 'scatter':
                ax.scatter(x, y, color=accent, s=100, edgecolors='white', linewidth=1.5)
            else:
                ax.plot(x, y, color=accent, linewidth=2.5, marker='o',
                        markersize=8, markerfacecolor=secondary, markeredgecolor='white')

            ax.grid(True, alpha=0.2, color=[c / 255.0 for c in theme['text']])
            plt.tight_layout()

            fd, path = tempfile.mkstemp(suffix='.png')
            os.close(fd)
            plt.savefig(path, dpi=150, bbox_inches='tight',
                        facecolor=fig.get_facecolor(), edgecolor='none')
            plt.close(fig)
            return path
        except Exception as e:
            if verbose:
                print(f"[PPTX] Chart generation failed: {e}")
            return None

    def _create_diagram_visual(self, output_path: Path, prompt: str,
                               theme: Dict[str, Any], verbose: bool) -> bool:
        """Create a diagram using matplotlib (enhanced box-and-arrow)."""
        try:
            import matplotlib
            matplotlib.use('Agg')
            import matplotlib.pyplot as plt
            import matplotlib.patches as mpatches
            from matplotlib.patches import FancyBboxPatch

            fig, ax = plt.subplots(figsize=(10, 5.5))
            bg_rgb = theme['bg']
            ax.set_facecolor([c / 255.0 for c in bg_rgb])
            fig.patch.set_facecolor([c / 255.0 for c in bg_rgb])

            # Extract keywords (capitalized words from prompt)
            words = re.findall(r'\b[A-Z][a-zA-Z]{2,20}\b', prompt)
            unique = []
            seen = set()
            for w in words:
                wl = w.lower()
                if wl not in seen and len(unique) < 8:
                    seen.add(wl)
                    unique.append(w)
            if not unique:
                unique = ["Input", "Process", "Output"]

            colors = [
                '#E69F00', '#56B4E9', '#009E73', '#F0E442',
                '#0072B2', '#D55E00', '#CC79A7', '#999999'
            ]
            text_color = '#333333' if theme['name'] == 'light' else '#CCCCCC'

            n = len(unique)
            cols = min(4, n)
            box_w, box_h = 2.0, 1.3
            spacing = 2.4
            start_x = (10 - (cols - 1) * spacing - box_w) / 2
            start_y = 6.5

            positions = []
            for i in range(n):
                col = i % cols
                row = i // cols
                x = start_x + col * spacing
                y = start_y - row * 2.5
                positions.append((x, y))

            for i, (word, (x, y)) in enumerate(zip(unique, positions)):
                color = colors[i % len(colors)]
                box = FancyBboxPatch(
                    (x, y), box_w, box_h,
                    boxstyle="round,pad=0.1",
                    facecolor=color, edgecolor='white',
                    linewidth=2, alpha=0.85
                )
                ax.add_patch(box)
                ax.text(x + box_w / 2, y + box_h / 2, word,
                        ha='center', va='center', fontsize=12,
                        fontweight='bold', color='white')

            # Arrows
            for i in range(n - 1):
                x1, y1 = positions[i]
                x2, y2 = positions[i + 1]
                if i % cols == cols - 1 and (i + 1) % cols == 0:
                    # Wrap to next row
                    ax.annotate('', xy=(x2 + box_w / 2, y2 + box_h),
                                xytext=(x1 + box_w / 2, y1),
                                arrowprops=dict(arrowstyle='->', lw=2, color='#666666'))
                else:
                    ax.annotate('', xy=(x2, y2 + box_h / 2),
                                xytext=(x1 + box_w, y1 + box_h / 2),
                                arrowprops=dict(arrowstyle='->', lw=2, color='#666666'))

            ax.set_xlim(-0.5, 10.5)
            ax.set_ylim(0, 8)
            ax.axis('off')
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight',
                        facecolor=fig.get_facecolor(), edgecolor='none')
            plt.close(fig)
            return True
        except Exception as e:
            if verbose:
                print(f"[PPTX] Diagram generation failed: {e}")
            return False

    # ═══════════════════════════════════════════════════════════════════
    # Prompt Parsing Utilities
    # ═══════════════════════════════════════════════════════════════════

    def _parse_prompt(self, prompt: str) -> Dict[str, Any]:
        """Extract title and bullets from prompt (enhanced)."""
        # Strip layout/theme tags
        clean = re.sub(r'\[(?:LAYOUT|THEME):\s*\w+\]', '', prompt, flags=re.IGNORECASE).strip()

        title = "Slide"
        bullets = []

        # Try to find title
        title_match = re.search(r'[Tt]itle[\s:]+["\']?([^"\'\n]+)["\']?', clean)
        if title_match:
            title = title_match.group(1).strip()
        else:
            first_line = clean.split('\n')[0].strip()
            if len(first_line) < 80:
                title = first_line

        # Extract bullets: numbered lists, dashes, key points
        for m in re.finditer(r'(?:^|\n)\s*(?:\d+[.):]|[-•*])\s+(.+)', clean):
            b = m.group(1).strip()
            if b and b not in bullets:
                bullets.append(b)

        # Key points section
        kp = re.search(r'[Kk]ey\s+[Pp]oints?[\s:]+(.+?)(?:\n\n|\Z)', clean, re.DOTALL)
        if kp:
            for line in kp.group(1).split('\n'):
                line = line.strip().lstrip('-•*').strip()
                if line and line not in bullets:
                    bullets.append(line)

        # Fallback: lines after first
        if not bullets:
            for line in clean.split('\n')[1:]:
                line = line.strip()
                if line and len(line) < 200:
                    bullets.append(line)

        return {"title": title, "bullets": bullets}

    def _extract_field(self, prompt: str, field_names: List[str]) -> str:
        """Extract a field like 'Title: value' from prompt."""
        for name in field_names:
            pattern = rf'{re.escape(name)}[\s:]+["\']?([^"\'\n]+)["\']?'
            m = re.search(pattern, prompt, re.IGNORECASE)
            if m:
                return m.group(1).strip()
        return ''

    def _parse_image_paths(self, prompt: str) -> List[str]:
        """Extract image paths from IMAGE: markers."""
        matches = re.findall(
            r'IMAGE:\s*(\S+\.(?:png|jpg|jpeg|gif|bmp|svg))',
            prompt, re.IGNORECASE
        )
        return matches

    def _parse_table_data(self, prompt: str) -> Optional[Dict[str, Any]]:
        """Parse TABLE: marker or markdown table."""
        # Find TABLE: section
        table_match = re.search(
            r'TABLE:\s*(.+?)(?:\n\n|\Z)',
            prompt, re.IGNORECASE | re.DOTALL
        )
        if not table_match:
            return None

        lines = [l.strip() for l in table_match.group(1).strip().split('\n') if l.strip()]
        if not lines:
            return None

        # Try pipe-delimited format: Header1 | Header2
        if '|' in lines[0]:
            headers = [h.strip() for h in lines[0].split('|')]
            rows = []
            for line in lines[1:]:
                if line.startswith('---') or line.startswith('==='):
                    continue
                row = [c.strip() for c in line.split('|')]
                if any(row):
                    rows.append(row)
            return {'headers': headers, 'rows': rows}

        # Try tab/comma format
        delim = '\t' if '\t' in lines[0] else ','
        headers = [h.strip() for h in lines[0].split(delim)]
        rows = []
        for line in lines[1:]:
            row = [c.strip() for c in line.split(delim)]
            if any(row):
                rows.append(row)
        return {'headers': headers, 'rows': rows}

    # ═══════════════════════════════════════════════════════════════════
    # Drawing Helpers
    # ═══════════════════════════════════════════════════════════════════

    def _add_bullet_column(self, slide, bullets, theme, left, top, width, height):
        """Add a column of bullets to a slide."""
        if not bullets:
            return
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet[:LAYOUT_MATRIX['two_column']['bullet_chars']]}"
            p.font.size = self._Pt(20)
            p.font.name = theme['font_body']
            p.font.color.rgb = self._rgb(theme['bullet'])
            p.space_after = self._Pt(14)

    def _add_placeholder_rect(self, slide, left, top, width, height, theme):
        """Add a placeholder rectangle with text."""
        rect = slide.shapes.add_shape(
            self._MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = self._RGBColor(230, 230, 230)
        rect.line.color.rgb = self._rgb(theme['accent'])
        rect.line.width = self._Pt(2)

        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "[Image Placeholder]"
        p.font.size = self._Pt(16)
        p.font.name = theme['font_body']
        p.font.color.rgb = self._RGBColor(150, 150, 150)
        p.alignment = self._PP_ALIGN.CENTER

    # ═══════════════════════════════════════════════════════════════════
    # QA Validation (inspired by Mck-ppt-design gate_check)
    # ═══════════════════════════════════════════════════════════════════

    def _qa_validate(self, prs) -> List[str]:
        """Validate slides for layout issues."""
        issues = []
        EMU_PER_INCH = 914400
        MARGIN = 0.2  # inch tolerance

        for si, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not hasattr(shape, 'top'):
                    continue
                top = (shape.top or 0) / EMU_PER_INCH
                left = (shape.left or 0) / EMU_PER_INCH
                w = (shape.width or 0) / EMU_PER_INCH
                h = (shape.height or 0) / EMU_PER_INCH
                bottom = top + h
                right = left + w

                if top < -MARGIN:
                    issues.append(f'Slide {si + 1}: shape above slide (top={top:.1f}")')
                if bottom > SLIDE_H_INCH + MARGIN:
                    issues.append(f'Slide {si + 1}: shape below slide (bottom={bottom:.1f}")')
                if left < -MARGIN:
                    issues.append(f'Slide {si + 1}: shape off left (left={left:.1f}")')
                if right > SLIDE_W_INCH + MARGIN:
                    issues.append(f'Slide {si + 1}: shape off right (right={right:.1f}")')
                if w < 0:
                    issues.append(f'Slide {si + 1}: negative width ({w:.2f}")')
                if h < 0:
                    issues.append(f'Slide {si + 1}: negative height ({h:.2f}")')

        return issues


if __name__ == "__main__":
    # Quick test
    import sys
    backend = PPTXBackend(Path(__file__).parent)
    result = backend.generate(
        "[LAYOUT: title] Machine Learning in Healthcare\nSubtitle: A Review\nAuthor: Dr. Smith\nDate: 2026-05-19",
        Path("/tmp/test_title.pptx")
    )
    print(result)
