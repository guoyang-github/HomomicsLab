#!/usr/bin/env python3
"""
Rearrange, duplicate, or reorder slides in PowerPoint presentations.
"""

import argparse
import sys
from pathlib import Path
from typing import List


def rearrange_slides(input_path: Path, output_path: Path, order: List[int]) -> None:
    """Rearrange slides according to the specified order."""
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(input_path)
    total_slides = len(prs.slides)

    # Validate indices
    for idx in order:
        if idx < 0 or idx >= total_slides:
            print(f"Error: Slide index {idx} out of range (0-{total_slides - 1})", file=sys.stderr)
            sys.exit(1)

    # python-pptx doesn't support direct slide reordering
    # We need to use a workaround with xml manipulation
    from copy import deepcopy

    # Get the slide XML elements
    slide_elements = list(prs.slides._sldIdLst)
    slide_layouts = [slide.slide_layout for slide in prs.slides]

    # Create a new presentation
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    # We need to rebuild the slide order
    # This is complex with python-pptx, so we'll use a simpler approach:
    # Save slides as images/rebuild or use a different strategy

    # Better approach: manipulate the XML directly
    sldIdLst = prs.part._element.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst")

    if sldIdLst is None:
        print("Error: Could not find slide ID list in presentation", file=sys.stderr)
        sys.exit(1)

    # Get all sldId elements in current order
    sldIds = list(sldIdLst)

    # Clear the list
    for child in list(sldIdLst):
        sldIdLst.remove(child)

    # Add back in new order (duplicates allowed)
    for idx in order:
        sldIdLst.append(deepcopy(sldIds[idx]))

    prs.save(output_path)
    print(f"Saved rearranged presentation to {output_path}")
    print(f"New order: {order}")
    print(f"Total slides: {len(order)}")


def main():
    parser = argparse.ArgumentParser(description="Rearrange slides in PowerPoint presentations")
    parser.add_argument("input", help="Input PowerPoint file")
    parser.add_argument("output", help="Output PowerPoint file")
    parser.add_argument("order", help="Comma-separated slide indices (0-based), duplicates allowed")
    args = parser.parse_args()

    try:
        order = [int(x.strip()) for x in args.order.split(",")]
    except ValueError:
        print("Error: --order must be comma-separated integers (e.g., '2,0,1')", file=sys.stderr)
        sys.exit(1)
    rearrange_slides(Path(args.input), Path(args.output), order)


if __name__ == "__main__":
    main()
